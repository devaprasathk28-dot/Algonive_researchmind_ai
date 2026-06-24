from pptx import Presentation

def generate_ppt_report(
    filename,
    title,
    summary
):

    prs = Presentation()

    slide_layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(
        slide_layout
    )

    slide.shapes.title.text = title

    slide.placeholders[1].text = summary

    prs.save(filename)

    return filename
